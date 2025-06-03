import streamlit as st
import uuid
import pandas as pd
import numpy as np
import io
from openai import OpenAI
from datetime import datetime
import altair as alt
import matplotlib.pyplot as plt
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Initialize the uploader key for true reset ---
if "uploader_key" not in st.session_state:
    st.session_state["uploader_key"] = str(uuid.uuid4())

# --- Load OpenAI API key
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

st.title("🧠 ClaimSight – Reimbursement Forecasting & Delay Analyzer")
st.image("claimsight_logo.png", width=200)

st.markdown(
    """
    <div style="background-color:#ffeaea; padding:15px; border:2px solid #ff3b3b; border-radius:7px;">
        <span style="color:#d90429; font-weight:bold; font-size:1.1em;">
            ⚠️ Please ensure your claims file does NOT include any patient names, dates of birth, or PHI.<br>
            Files should be de-identified before upload.<br>
            <b>Any PHI is ignored and never retained or stored.</b>
        </span>
    </div>
    """,
    unsafe_allow_html=True,
)

# 👇 Add instructions for sample file and CSV flexibility here!
st.markdown("""
👋 **Want to see how ClaimSight works?**
- [⬇️ Download a Sample Claims File](https://drive.google.com/file/d/1mc03haJ6znwXFBuqrNTK-aOxtQ_OYn1g/view?usp=sharing)
- Or upload your own de-identified claims file—ClaimSight works with any CSV format!
""")

# ---- USE THE SESSION KEY ----
uploaded_file = st.file_uploader(
    "📤 Upload your claim CSV (any format/order/labels!)",
    type=["csv"],
    key=st.session_state["uploader_key"]  
)
st.markdown("_Your data is processed securely and never stored. All uploads are de-identified._")
st.markdown(
    '<a href="https://claimsight.carrd.co/" target="_blank" style="text-decoration:none;">'
    '🌐 <b>Learn more at claimsight.carrd.co</b>'
    '</a>',
    unsafe_allow_html=True
)

if uploaded_file:
    # --- Flexible import (column aliasing)
    field_map = {
        "Claim_ID": ["Claim_ID", "ClaimId", "Claim Number", "ClaimNumber", "ID"],
        "Claim_Date": ["Claim_Date", "Service Date", "Date of Service", "DOS"],
        "Payment_Date": ["Payment_Date", "Date Paid", "Paid_Date", "Paid Date"],
        "Payer": ["Payer", "Insurer", "Insurance", "Insurance Name"],
        "CPT_Code": ["CPT_Code", "CPT", "Procedure Code", "Procedure"],
        "Billed_Amount": ["Billed_Amount", "Amount Billed", "Charge Amount", "Charge"],
        "Paid_Amount": ["Paid_Amount", "Amount Paid", "Paid", "Payment"],
        "Status": ["Status", "Claim Status", "Claim_Status", "Payment Status", "State"]
    }
    df_raw = pd.read_csv(uploaded_file)
    df = pd.DataFrame()
    for std_col, aliases in field_map.items():
        for alt_col in aliases:
            if alt_col in df_raw.columns:
                df[std_col] = df_raw[alt_col]
                break
        else:
            df[std_col] = np.nan

    # Parse dates, fix types
    df['Claim_Date'] = pd.to_datetime(df['Claim_Date'], errors='coerce')
    df['Payment_Date'] = pd.to_datetime(df['Payment_Date'], errors='coerce')
    for col in ['Billed_Amount', 'Paid_Amount']:
        df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)
    df['Status'] = df['Status'].astype(str)

    st.subheader("📊 Claim Summary")
    st.dataframe(df)

    # --- KPIs
    total_billed = df['Billed_Amount'].sum()
    total_paid = df['Paid_Amount'].sum()
    pending_count = (df['Status'].str.lower() == 'pending').sum()
    st.metric("Total Claims", len(df))
    st.metric("💸 Total Billed", f"${total_billed:,.2f}")
    st.metric("💰 Total Paid", f"${total_paid:,.2f}")
    st.metric("⏳ Pending Claims", pending_count)

    # --- Data Quality Alerts (Restored!)
    today = pd.Timestamp.today().normalize()
    dq_issues = df[
        (df['Claim_Date'].isna()) |
        (df['Billed_Amount'] <= 0) |
        (df['Paid_Amount'] < 0) |
        (df['Status'].isin(['nan', '', 'None', 'null', 'NULL']))
    ].copy()
    dq_issues['Days Late'] = (today - dq_issues['Claim_Date']).dt.days
    if not dq_issues.empty:
        st.subheader("⚠️ Data Quality Alerts")
        st.dataframe(dq_issues, use_container_width=True)

    # --- Aging Buckets (FIXED ORDER)
    def get_aging_bucket(row):
        if pd.isnull(row['Payment_Date']):
            days = (today - row['Claim_Date']).days
        else:
            days = (row['Payment_Date'] - row['Claim_Date']).days
        if days <= 30:
            return '0-30'
        elif days <= 60:
            return '31-60'
        elif days <= 90:
            return '61-90'
        elif days <= 120:
            return '91-120'
        else:
            return '120+'

    df['Aging Bucket'] = df.apply(get_aging_bucket, axis=1)
    bucket_order = ["0-30", "31-60", "61-90", "91-120", "120+"]
    df['Aging Bucket'] = pd.Categorical(df['Aging Bucket'], categories=bucket_order, ordered=True)
    bucket_grouped = df.groupby('Aging Bucket').agg(
        Count=('Claim_ID', 'count'),
        Total_Amount=('Billed_Amount', 'sum'),
        Paid=('Paid_Amount', 'sum')
    ).reset_index()
    bucket_grouped['Outstanding'] = bucket_grouped['Total_Amount'] - bucket_grouped['Paid']
    bucket_summary = bucket_grouped.sort_values('Aging Bucket')
    st.subheader("📆 Aging Buckets")
    st.dataframe(bucket_summary, use_container_width=True)
    st.bar_chart(bucket_summary.set_index("Aging Bucket")["Total_Amount"])

    # --- Payer Mix by Dollar
    st.subheader("🏦 Payer Mix by Dollar")
    payer_mix = df.groupby('Payer')['Billed_Amount'].sum().sort_values(ascending=False)
    st.dataframe(payer_mix.rename("Billed Amount ($)"))
    fig1, ax1 = plt.subplots()
    payer_mix.plot.pie(autopct='%1.1f%%', ylabel='', ax=ax1)
    ax1.set_title("Payer Mix by Dollar")
    st.pyplot(fig1)

    # --- Denials & Write-Offs
    st.subheader("🚫 Denials & Write-offs")
    denial_rate = (df['Status'].str.lower() == 'denied').mean()
    writeoff_rate = (df['Status'].str.lower().isin(['write-off', 'writeoff', 'written off'])).mean()
    st.metric("❌ Denial Rate", f"{denial_rate:.1%}")
    st.metric("📝 Write-off Rate", f"{writeoff_rate:.1%}")

    denial_by_payer = df.groupby('Payer').apply(lambda x: (x['Status'].str.lower() == 'denied').mean())
    writeoff_by_payer = df.groupby('Payer').apply(lambda x: (x['Status'].str.lower().isin(['write-off', 'writeoff', 'written off'])).mean())
    st.bar_chart(denial_by_payer.rename("Denial Rate by Payer"))
    st.bar_chart(writeoff_by_payer.rename("Write-off Rate by Payer"))

    # --- Largest Individual Outstanding Claims
    st.subheader("💰 Largest Individual Outstanding Claims")
    df['Outstanding'] = df['Billed_Amount'] - df['Paid_Amount']
    top_outstanding = df.sort_values("Outstanding", ascending=False).head(10)
    st.dataframe(top_outstanding[['Claim_ID', 'Payer', 'Claim_Date', 'Billed_Amount', 'Paid_Amount', 'Outstanding']])

    # --- Collection Rates
    gross_collection_rate = total_paid / total_billed if total_billed > 0 else 0
    net_collection_rate = total_paid / (total_billed - df[df['Status'].str.lower().isin(['write-off', 'writeoff', 'written off'])]['Billed_Amount'].sum()) if total_billed > 0 else 0

    gross_collection_rate_by_payer = df.groupby('Payer').apply(
        lambda x: x['Paid_Amount'].sum() / x['Billed_Amount'].sum() if x['Billed_Amount'].sum() > 0 else 0)
    net_collection_rate_by_payer = df.groupby('Payer').apply(
        lambda x: x['Paid_Amount'].sum() / (x['Billed_Amount'].sum() - x[x['Status'].str.lower().isin(['write-off', 'writeoff', 'written off'])]['Billed_Amount'].sum())
        if (x['Billed_Amount'].sum() - x[x['Status'].str.lower().isin(['write-off', 'writeoff', 'written off'])]['Billed_Amount'].sum()) > 0 else 0)

    st.markdown("💳 **Collection Rates**")
    col_gross, col_net = st.columns(2)
    with col_gross:
        st.markdown("**Gross Collection Rate**")
        st.markdown(f"<h2>{gross_collection_rate:.1%}</h2>", unsafe_allow_html=True)
    with col_net:
        st.markdown("**Net Collection Rate**")
        st.markdown(f"<h2>{net_collection_rate:.1%}</h2>", unsafe_allow_html=True)
    st.markdown("#### Gross Collection Rate by Payer")
    st.bar_chart(gross_collection_rate_by_payer)
    st.markdown("#### Net Collection Rate by Payer")
    st.bar_chart(net_collection_rate_by_payer)

    # --- Avg Payer Delay (Days)
    st.subheader("📉 Avg Payer Delay (Days)")
    paid_df = df[df['Status'].str.lower() == "paid"].copy()
    paid_df["Delay"] = (paid_df["Payment_Date"] - paid_df["Claim_Date"]).dt.days
    delay_by_payer = paid_df.groupby("Payer")["Delay"].mean().round(1)
    st.dataframe(delay_by_payer.rename("Avg Delay (days)"))
    st.bar_chart(delay_by_payer)

    # --- Claim Status Breakdown
    st.subheader("🥧 Claim Status Breakdown")
    status_counts = df['Status'].value_counts().reset_index()
    status_counts.columns = ['Status', 'Count']
    st.dataframe(status_counts)
    fig2, ax2 = plt.subplots()
    status_counts.set_index('Status')['Count'].plot.pie(autopct='%1.1f%%', ylabel='', title='Claim Status Breakdown', ax=ax2)
    st.pyplot(fig2)

    # --- Weekly Payment Timeline
    st.subheader("📆 Weekly Payment Timeline Forecast (Paid + Forecasted)")
    df['Claim_Date'] = pd.to_datetime(df['Claim_Date'])
    df['Payment_Date'] = pd.to_datetime(df['Payment_Date'])
    paid_df = df[df['Status'].str.lower() == 'paid'].copy()
    pending_df = df[df['Status'].str.lower() == 'pending'].copy()
    if 'Delay' not in paid_df.columns:
        paid_df["Delay"] = (paid_df["Payment_Date"] - paid_df["Claim_Date"]).dt.days
    paid_df['Week'] = paid_df['Payment_Date'].dt.to_period('W').apply(lambda r: r.start_time)
    weekly_paid = paid_df.groupby('Week')['Paid_Amount'].sum()
    if not pending_df.empty:
        avg_delay = paid_df.groupby("Payer")["Delay"].mean()
        pending_df['Expected_Days'] = pending_df['Payer'].map(avg_delay).fillna(14).astype(int)
        pending_df['Expected_Payment_Date'] = pending_df['Claim_Date'] + pd.to_timedelta(pending_df['Expected_Days'], unit='D')
        pending_df['Week'] = pending_df['Expected_Payment_Date'].dt.to_period('W').apply(lambda r: r.start_time)
        weekly_forecast = pending_df.groupby('Week')['Billed_Amount'].sum()
    else:
        weekly_forecast = pd.Series(dtype=float)
    all_weeks = pd.date_range(
        start=df['Claim_Date'].min().to_period('W').start_time,
        end=max(
            paid_df['Payment_Date'].max() if not paid_df.empty else df['Claim_Date'].max(),
            pending_df['Expected_Payment_Date'].max() if not pending_df.empty else df['Claim_Date'].max()
        ),
        freq='W-MON'
    )
    full_timeline = pd.DataFrame(index=all_weeks)
    full_timeline['Paid'] = weekly_paid
    full_timeline['Forecast'] = weekly_forecast
    full_timeline = full_timeline.fillna(0).round(2)
    full_timeline['Week'] = full_timeline.index
    full_timeline = full_timeline.reset_index(drop=True)
    full_long = full_timeline.melt(id_vars=['Week'], value_vars=['Paid', 'Forecast'],
                                   var_name='Type', value_name='Amount')
    area_chart = alt.Chart(full_long).mark_area(opacity=0.7).encode(
        x=alt.X('Week:T', title='Week'),
        y=alt.Y('Amount:Q', stack='zero', title='Amount ($)'),
        color=alt.Color('Type:N', scale=alt.Scale(scheme='tableau10'), title='Type'),
        tooltip=['Week:T', 'Type:N', 'Amount:Q']
    ).properties(width=700, height=400)
    st.altair_chart(area_chart, use_container_width=True)

    # --- GPT Section ---
    st.subheader("🧠 Generate Expert Summary")
    if st.button("Generate GPT-4 Summary"):
        with st.spinner("Analyzing with GPT..."):
            prompt = f"""
            You are a healthcare finance analyst. Review the following claim dataset:
            - Identify payers with delayed or pending payments
            - Highlight CPT codes associated with high delays or denials
            - Surface data quality errors or coding issues
            - Anticipate cash flow risk, outstanding claims, and aging problems
            - Suggest specific actions to improve claim outcomes and revenue

            Dataset:
            {df.to_csv(index=False)}
            """
            response = client.chat.completions.create(
                model="gpt-4-turbo",
                messages=[
                    {"role": "system", "content": "You are an expert in healthcare reimbursement analytics."},
                    {"role": "user", "content": prompt}
                ]
            )
            st.session_state['summary'] = response.choices[0].message.content
            st.success("GPT Summary generated.")

    if 'summary' in st.session_state:
        st.markdown(st.session_state['summary'])

    # --- Download Section ---
    st.subheader("📥 Download Report Files")
    col1, col2, col3 = st.columns(3)
    with col1:
        output = io.BytesIO()
        export_timeline = full_timeline.copy()
        export_timeline['Total'] = export_timeline['Paid'] + export_timeline['Forecast']
        export_timeline['Week'] = pd.to_datetime(export_timeline['Week']).dt.strftime('%Y-%m-%d')
        df_export = df.copy()
        df_export['Claim_Date'] = pd.to_datetime(df_export['Claim_Date']).dt.strftime('%Y-%m-%d')
        df_export['Payment_Date'] = pd.to_datetime(df_export['Payment_Date'], errors='coerce').dt.strftime('%Y-%m-%d')
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            df_export.to_excel(writer, sheet_name='Original Claims', index=False)
            export_timeline.to_excel(writer, sheet_name='Weekly Forecast', index=False)
            delay_by_payer.to_frame(name='Avg Delay').to_excel(writer, sheet_name='Avg Delay by Payer')
            status_counts.to_excel(writer, sheet_name='Claim Status Breakdown', index=False)
            bucket_summary.to_excel(writer, sheet_name='Aging Buckets', index=False)
            payer_mix.to_frame(name='Billed Amount ($)').to_excel(writer, sheet_name='Payer Mix', index=True)
            denial_by_payer.to_frame(name='Denial Rate').to_excel(writer, sheet_name='Denial Rate by Payer')
            writeoff_by_payer.to_frame(name='Write-off Rate').to_excel(writer, sheet_name='Write-off Rate by Payer')
            gross_collection_rate_by_payer.to_frame(name='Gross Collection Rate').to_excel(writer, sheet_name='Gross Collection by Payer')
            net_collection_rate_by_payer.to_frame(name='Net Collection Rate').to_excel(writer, sheet_name='Net Collection by Payer')
            top_outstanding.to_excel(writer, sheet_name='Largest Outstanding Claims', index=False)
            if 'summary' in st.session_state:
                pd.DataFrame({'GPT_Summary': [st.session_state['summary']]}).to_excel(writer, sheet_name='AI Summary', index=False)
        st.download_button(
            "Download Excel",
            data=output.getvalue(),
            file_name="ClaimSight_Report.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    with col2:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=LETTER)
        styles = getSampleStyleSheet()
        story = []
        story.append(Paragraph("ClaimSight Summary Report", styles["Title"]))
        story.append(Spacer(1, 12))
        story.append(Paragraph(f"Total Claims: {len(df)}", styles["Normal"]))
        story.append(Paragraph(f"Total Billed: ${total_billed:,.2f}", styles["Normal"]))
        story.append(Paragraph(f"Total Paid: ${total_paid:,.2f}", styles["Normal"]))
        story.append(Spacer(1, 12))
        if 'summary' in st.session_state:
            story.append(Paragraph("GPT Summary", styles["Heading2"]))
            for line in st.session_state['summary'].splitlines():
                story.append(Paragraph(line.strip(), styles["Normal"]))
            story.append(PageBreak())
        def add_chart(title, plot_func):
            story.append(Spacer(1, 12))
            story.append(Paragraph(title, styles["Heading2"]))
            buf = io.BytesIO()
            plt.figure(figsize=(7, 3))
            plot_func()
            plt.tight_layout()
            plt.savefig(buf, format='png')
            plt.close()
            buf.seek(0)
            story.append(Image(buf, width=500, height=220))
            story.append(PageBreak())
        add_chart("Aging Buckets", lambda: bucket_summary.set_index("Aging Bucket")["Total_Amount"].plot(kind='bar'))
        add_chart("Payer Mix by Dollar", lambda: payer_mix.plot.pie(autopct='%1.1f%%', ylabel=''))
        add_chart("Denial Rate by Payer", lambda: denial_by_payer.plot(kind='bar'))
        add_chart("Write-off Rate by Payer", lambda: writeoff_by_payer.plot(kind='bar'))
        add_chart("Largest Individual Outstanding Claims", lambda: top_outstanding.set_index('Claim_ID')['Outstanding'].plot(kind='bar'))
        add_chart("Gross Collection Rate by Payer", lambda: gross_collection_rate_by_payer.plot(kind='bar', color='skyblue'))
        add_chart("Net Collection Rate by Payer", lambda: net_collection_rate_by_payer.plot(kind='bar', color='skyblue'))
        add_chart("Avg Payer Delay", lambda: delay_by_payer.sort_values().plot(kind='barh', color='#1f77b4'))
        add_chart("Claim Status Breakdown", lambda: status_counts.set_index('Status')['Count'].plot.pie(autopct='%1.1f%%', ylabel=''))
        add_chart("Weekly Payment Timeline", lambda: plt.stackplot(pd.to_datetime(full_timeline['Week']), full_timeline['Paid'], full_timeline['Forecast'], labels=["Paid", "Forecast"], colors=['#1f77b4', '#ff7f0e']))
        doc.build(story)
        buffer.seek(0)
        st.download_button(
            "Download PDF",
            data=buffer.getvalue(),
            file_name="ClaimSight_Report.pdf",
            mime="application/pdf"
        )
    with col3:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ClaimSight Summary Report"
        slide.placeholders[1].text = f"Total Claims: {len(df)}\nTotal Billed: ${total_billed:,.2f}\nTotal Paid: ${total_paid:,.2f}"
        if 'summary' in st.session_state:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "GPT Summary"
            tf = slide.placeholders[1].text_frame
            for line in st.session_state['summary'].splitlines():
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
        def add_ppt_chart(title, plot_func):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
            title_tf = title_box.text_frame
            title_tf.text = title
            buf = io.BytesIO()
            plt.figure(figsize=(7, 3))
            plot_func()
            plt.tight_layout()
            plt.savefig(buf, format='png')
            plt.close()
            buf.seek(0)
            slide.shapes.add_picture(buf, Inches(1), Inches(1.2), height=Inches(4.5))
        add_ppt_chart("Aging Buckets", lambda: bucket_summary.set_index("Aging Bucket")["Total_Amount"].plot(kind='bar'))
        add_ppt_chart("Payer Mix by Dollar", lambda: payer_mix.plot.pie(autopct='%1.1f%%', ylabel=''))
        add_ppt_chart("Denial Rate by Payer", lambda: denial_by_payer.plot(kind='bar'))
        add_ppt_chart("Write-off Rate by Payer", lambda: writeoff_by_payer.plot(kind='bar'))
        add_ppt_chart("Largest Individual Outstanding Claims", lambda: top_outstanding.set_index('Claim_ID')['Outstanding'].plot(kind='bar'))
        add_ppt_chart("Gross Collection Rate by Payer", lambda: gross_collection_rate_by_payer.plot(kind='bar', color='skyblue'))
        add_ppt_chart("Net Collection Rate by Payer", lambda: net_collection_rate_by_payer.plot(kind='bar', color='skyblue'))
        add_ppt_chart("Avg Payer Delay", lambda: delay_by_payer.sort_values().plot(kind='barh', color='#1f77b4'))
        add_ppt_chart("Claim Status Breakdown", lambda: status_counts.set_index('Status')['Count'].plot.pie(autopct='%1.1f%%', ylabel=''))
        add_ppt_chart("Weekly Payment Timeline", lambda: plt.stackplot(pd.to_datetime(full_timeline['Week']), full_timeline['Paid'], full_timeline['Forecast'], labels=["Paid", "Forecast"], colors=['#1f77b4', '#ff7f0e']))
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        st.download_button(
            "Download PowerPoint",
            data=output.getvalue(),
            file_name="ClaimSight_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
# --- RESET BUTTON (ALWAYS VISIBLE AT BOTTOM OF APP) ---
st.markdown("---")
if st.button("🔄 Reset / Start Over"):
    st.session_state.clear()  # Clear everything
    st.session_state["uploader_key"] = str(uuid.uuid4())  # New key for uploader
    st.rerun()  # Full rerun, now uploader is also reset!

st.markdown("---")
st.markdown(
    '💡 **Have feedback or ideas? [Let us know!](https://docs.google.com/forms/d/e/1FAIpQLSfm7kcmuA9pg_oenYrFQ62exJsvYcNHKV9Zb7b8RzaUOz_KHg/viewform?usp=sharing&ouid=104679907223026161047)**',
    unsafe_allow_html=True
)

st.markdown("---")
st.markdown(
    '<a href="https://claimsight.carrd.co/" target="_blank" style="text-decoration:none;">'
    '🌐 <b>Learn more at claimsight.carrd.co</b>'
    '</a>',
    unsafe_allow_html=True
)